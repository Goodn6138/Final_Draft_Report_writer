import os
from pptx import Presentation
from docx import Document
from openai import OpenAI

# Configure the LLM client. This uses the `openai` package but configured for a router endpoint (HuggingFace)
# Set environment variables: HF_ROUTER_API_KEY or OPENAI_API_KEY depending on provider

def get_llm_client():
    api_key = os.getenv('HF_ROUTER_API_KEY') or os.getenv('OPENAI_API_KEY')
    base_url = os.getenv('LLM_BASE_URL', 'https://router.huggingface.co/v1')
    if not api_key:
        raise ValueError('Set HF_ROUTER_API_KEY or OPENAI_API_KEY as environment variable')
    client = OpenAI(base_url=base_url, api_key=api_key)
    return client


def extract_text_from_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    texts.append(shape.text)
            except Exception:
                # ignore shapes that don't have text
                pass
    return '\n'.join(texts)


def generate_section_text(client, system_prompt, user_prompt, temperature=0.3, max_tokens=800):
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt}
    ]
    completion = client.chat.completions.create(
        model=os.getenv('LLM_MODEL', 'openai/gpt-oss-120b:cerebras'),
        messages=messages,
        temperature=temperature,
        max_tokens=max_tokens
    )
    return completion.choices[0].message.content


def generate_sections_from_template(template, proposal_text, notes):
    client = get_llm_client()
    system_prompt = (
        "You are an academic writing assistant. Use the provided proposal text and the user's notes "
        "to write high-quality academic sections. Respect the section titles and lengths from the template."
    )

    sections = {}
    for section in template.get('sections', []):
        title = section.get('title')
        guidance = section.get('guidance', '')
        length = section.get('words', 250)

        user_prompt = (
            f"Write the section titled \"{title}\". Target length: {length} words.\n\n"
            f"Guidance: {guidance}\n\n"
            f"Proposal Text:\n{proposal_text}\n\n"
            f"Additional Notes:\n{notes}\n\n"
            "If information is missing, write a clear placeholder and suggest what the student should add."
        )

        try:
            content = generate_section_text(client, system_prompt, user_prompt)
        except Exception as e:
            content = f"[ERROR generating section: {e}]"
        sections[title] = content
    return sections


def save_sections_to_docx(sections, out_path):
    doc = Document()
    for title, content in sections.items():
        doc.add_heading(title, level=1)
        doc.add_paragraph(content)
    doc.save(out_path)
    return out_path
